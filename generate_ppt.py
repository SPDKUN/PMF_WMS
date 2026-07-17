#!/usr/bin/env python3
"""生成答辩PPT：预制菜WMS智能仓储管理系统"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 颜色主题：仓储物流蓝 ──
BG_DARK   = RGBColor(0x0F, 0x1A, 0x2E)   # 深蓝背景
BG_CARD   = RGBColor(0x16, 0x24, 0x3E)   # 卡片深色
BLUE_MAIN = RGBColor(0x3B, 0x82, 0xF6)   # 主蓝色
BLUE_LIGHT= RGBColor(0x60, 0xA5, 0xFA)   # 浅蓝
TEAL      = RGBColor(0x14, 0xB8, 0xA6)   # 青绿
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0x94, 0xA3, 0xB8)
DARK_TEXT  = RGBColor(0xE2, 0xE8, 0xF0)
GOLD      = RGBColor(0xFB, 0xBF, 0x24)
RED_ACCENT= RGBColor(0xEF, 0x44, 0x44)
GREEN_OK  = RGBColor(0x22, 0xC5, 0x5E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── 工具函数 ──
def add_bg(slide, color=BG_DARK):
    """填充纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = 0
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=16, default_color=DARK_TEXT):
    """添加多行文本，lines=[(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text = line[0]
        size = line[1] if len(line) > 1 else default_size
        color = line[2] if len(line) > 2 else default_color
        bold = line[3] if len(line) > 3 else False
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Microsoft YaHei'
    return txBox

def add_accent_line(slide, left, top, width, color=BLUE_MAIN):
    """添加装饰线条"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_image_placeholder(slide, left, top, width, height):
    """添加图片占位框（虚线边框+文字提示）"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x2E, 0x4A)
    shape.line.color.rgb = BLUE_MAIN
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 2  # dash
    # 添加提示文字
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📷  插入图片"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_tag(slide, left, top, text, bg_color, text_color=WHITE):
    """添加小标签"""
    shape = add_rect(slide, left, top, Inches(1.8), Inches(0.35), bg_color)
    shape.fill.solid()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    return shape

# ──────────────────────────────────────────────
# 第1页：封面
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BG_DARK)

# 顶部装饰条
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), BLUE_MAIN)

# 左侧装饰竖线
add_rect(slide, Inches(1.5), Inches(1.8), Pt(4), Inches(2.8), TEAL)

# 主标题
add_text_box(slide, Inches(2.0), Inches(1.8), Inches(10), Inches(1.2),
             "预制菜 WMS 智能仓储管理系统", font_size=48, color=WHITE, bold=True)

# 副标题
add_text_box(slide, Inches(2.0), Inches(3.0), Inches(10), Inches(0.8),
             "Prepared Food Warehouse Management System", font_size=22, color=BLUE_LIGHT)

# 装饰线
add_accent_line(slide, Inches(2.0), Inches(3.8), Inches(3), TEAL)

# 信息区
add_multi_text(slide, Inches(2.0), Inches(4.2), Inches(8), Inches(1.8), [
    ("答辩时间：2026年7月19日", 18, LIGHT_GRAY),
    ("汇报人：彭健坤    王佳宁    彭双华", 20, WHITE, True),
    ("第24组", 16, BLUE_LIGHT),
])

# 右下角装饰
add_rect(slide, Inches(11), Inches(6.5), Inches(2.5), Inches(1.2), RGBColor(0x1A, 0x30, 0x50))

# ──────────────────────────────────────────────
# 第2页：目录
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.7),
             "目  录", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), TEAL)

toc_items = [
    ("01", "项目概述", "背景、目标与系统定位"),
    ("02", "技术架构", "前后端技术栈与系统架构"),
    ("03", "核心业务流程", "出入库、质检、盘点的完整闭环"),
    ("04", "AI 智能助手", "DeepSeek 大模型 + 真实数据注入"),
    ("05", "数据可视化看板", "ECharts 图表实时监控"),
    ("06", "微信小程序", "跨平台移动端适配"),
    ("07", "项目亮点总结", "三大亮点与创新点"),
    ("08", "系统演示", "现场展示 + Q&A"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.6) + Inches(0.7) * i
    # 编号
    add_text_box(slide, Inches(0.8), y, Inches(0.6), Inches(0.5),
                 num, font_size=28, color=BLUE_MAIN, bold=True)
    # 竖线
    add_rect(slide, Inches(1.5), y + Inches(0.05), Pt(3), Inches(0.4),
             BLUE_MAIN if i < 7 else TEAL)
    # 标题
    add_text_box(slide, Inches(1.8), y, Inches(3), Inches(0.3),
                 title, font_size=20, color=WHITE, bold=True)
    # 描述
    add_text_box(slide, Inches(1.8), y + Inches(0.3), Inches(5), Inches(0.3),
                 desc, font_size=13, color=LIGHT_GRAY)

# 右侧图片占位
add_image_placeholder(slide, Inches(8.0), Inches(1.6), Inches(4.5), Inches(5.2))

# ──────────────────────────────────────────────
# 第3页：项目概述
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.7),
             "01  项目概述", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), TEAL)

# 左边内容区
add_multi_text(slide, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.0), [
    ("📋  项目定位", 22, WHITE, True),
    ("为预制菜行业量身打造的仓库管理系统（WMS），覆盖从货物入库、存储、", 15, DARK_TEXT),
    ("出库到质检、盘点的全生命周期管理。", 15, DARK_TEXT),
    ("", 8, WHITE),
    ("🎯  解决的核心问题", 22, WHITE, True),
    ("• 传统纸质记录效率低、易出错 → 数字化管理", 15, DARK_TEXT),
    ("• 库存数据不透明、难追溯 → 实时库存 + 操作日志", 15, DARK_TEXT),
    ("• 多角色协作混乱 → 主管/员工/质检员三级权限", 15, DARK_TEXT),
    ("• 缺乏数据支撑决策 → 可视化数据看板", 15, DARK_TEXT),
    ("", 8, WHITE),
    ("📊  系统规模", 22, WHITE, True),
    ("• 22 张数据库表 | 8 个前端页面 | 6 个小程序页面", 15, DARK_TEXT),
    ("• 支持 Web 端 + 微信小程序双平台", 15, DARK_TEXT),
    ("• 完整覆盖 入库→存储→出库→质检→盘点 全流程", 15, DARK_TEXT),
])

# 右侧图片占位
add_image_placeholder(slide, Inches(7.8), Inches(1.5), Inches(4.8), Inches(5.2))

# ──────────────────────────────────────────────
# 第4页：技术架构
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.7),
             "02  技术架构", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), TEAL)

# 架构层次
layers = [
    ("🖥  前端层", "Vue 3 + Element Plus + ECharts + Axios", BLUE_MAIN,
     "Wide-screen (1920×1080), 暗色/亮色主题切换, AES密码加密传输"),
    ("🔗  通信层", "HTTP RESTful API + JWT 令牌认证", TEAL,
     "Axios 拦截器自动附加 Token, Vite 代理解决跨域, 统一 ResultData 响应格式"),
    ("⚙  后端层", "Spring Boot 3 + MyBatis + Maven", BLUE_LIGHT,
     "三层架构(Controller→Service→Mapper), 事务管理(@Transactional), 全局异常处理"),
    ("💾  数据层", "MySQL 8.0 + Redis 6", BLUE_MAIN,
     "22张业务表, Redis 缓存用户会话, MyBatis 动态 SQL"),
    ("🤖  AI 层", "DeepSeek 大模型 (API 调用)", TEAL,
     "OkHttp 发送请求, 系统提示词约束, 真实数据库数据注入上下文"),
]

for i, (title, tech, accent, desc) in enumerate(layers):
    y = Inches(1.5) + Inches(1.05) * i
    add_rect(slide, Inches(0.8), y, Inches(6.5), Inches(0.9), BG_CARD)
    add_rect(slide, Inches(0.8), y, Pt(4), Inches(0.9), accent)
    add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(5.5), Inches(0.35),
                 title, font_size=17, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.38), Inches(3.5), Inches(0.25),
                 tech, font_size=13, color=accent)
    add_text_box(slide, Inches(1.2), y + Inches(0.6), Inches(5.5), Inches(0.25),
                 desc, font_size=11, color=LIGHT_GRAY)

# 右侧技术栈总览卡片
add_rect(slide, Inches(7.8), Inches(1.5), Inches(4.8), Inches(2.8), BG_CARD)
add_multi_text(slide, Inches(8.2), Inches(1.7), Inches(4.0), Inches(2.4), [
    ("技术栈一览", 18, WHITE, True),
    ("", 5, WHITE),
    ("前端：Vue 3 · Element Plus · ECharts", 14, BLUE_LIGHT),
    ("后端：Spring Boot 3 · MyBatis", 14, BLUE_LIGHT),
    ("数据库：MySQL · Redis", 14, BLUE_LIGHT),
    ("安全：JWT · BCrypt · AES-256-CBC", 14, TEAL),
    ("AI：DeepSeek · OkHttp", 14, TEAL),
    ("工具：EasyExcel · Vite · Maven", 14, LIGHT_GRAY),
])

# 右下图片占位
add_image_placeholder(slide, Inches(7.8), Inches(4.6), Inches(4.8), Inches(2.3))

# ──────────────────────────────────────────────
# 第5页：核心业务流程（亮点1）
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
             "03  核心业务流程 —— 完整闭环", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), TEAL)

# 流程步骤
flow_steps = [
    ("① 批次创建", "生产完成后创建批次\n记录生产日期、保质期", BLUE_MAIN),
    ("② 入库操作", "选择批次→指定库位\n9张表联动更新", TEAL),
    ("③ 库存管理", "实时库位占用状态\n先进先出(FIFO)", BLUE_LIGHT),
    ("④ 质检管理", "来料检/日常抽检\n合格放行·不合格报废", GOLD),
    ("⑤ 出库操作", "选择批次→指定库位\n库存扣减·日志记录", BLUE_MAIN),
    ("⑥ 报表看板", "数据可视化总览\n出入库趋势·温湿度", TEAL),
]

for i, (title, desc, accent) in enumerate(flow_steps):
    x = Inches(0.5) + Inches(2.1) * i
    y = Inches(1.6)
    add_rect(slide, x, y, Inches(1.9), Inches(2.2), BG_CARD)
    add_rect(slide, x, y, Inches(1.9), Pt(4), accent)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.15), Inches(1.7), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.65), Inches(1.7), Inches(1.3),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    # 箭头（除最后一个）
    if i < len(flow_steps) - 1:
        add_text_box(slide, x + Inches(1.9), y + Inches(0.8), Inches(0.25), Inches(0.4),
                     "→", font_size=24, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

# 下半部分：亮点说明
add_multi_text(slide, Inches(0.8), Inches(4.1), Inches(12), Inches(2.8), [
    ("✨  亮点：一次入库操作涉及 9 张数据库表的联动更新，使用 @Transactional 保证事务一致性", 16, GOLD, True),
    ("", 5, WHITE),
    ("涉及表：inbound_order_head → inbound_order_detail → work_task → operation_log → inventory →", 13, LIGHT_GRAY),
    ("         location → zone → warehouse → goods → batch（10张表全部通过事务保护，任何一步失败即回滚）", 13, LIGHT_GRAY),
    ("", 5, WHITE),
    ("💡  三层权限架构：主管（全部权限）→ 仓库员工（出入库执行）→ 质检员（质量检查），前端路由守卫 + 后端接口鉴权", 14, DARK_TEXT),
])

# 底部图片占位
add_image_placeholder(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.2))

# ──────────────────────────────────────────────
# 第6页：AI智能助手（亮点2）
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
             "04  AI 智能助手 —— DeepSeek 大模型接入", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), TEAL)

# 左边：实现流程
add_multi_text(slide, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.2), [
    ("🤖  实现原理", 22, WHITE, True),
    ("", 6, WHITE),
    ("1️⃣  用户输入问题 → 前端将问题+历史对话发送到后端 /ai/chat", 15, DARK_TEXT),
    ("2️⃣  后端分析关键词（库存/货物/仓库/批次…）→ 查询真实数据库", 15, DARK_TEXT),
    ("3️⃣  组装消息：系统提示词 + 真实数据上下文 + 历史对话 + 当前问题", 15, DARK_TEXT),
    ("4️⃣  OkHttp 发送到 DeepSeek API → 解析回复 → 返回前端", 15, DARK_TEXT),
    ("", 10, WHITE),
    ("🔑  核心技术点", 22, WHITE, True),
    ("", 6, WHITE),
    ("• 系统提示词工程：定义AI人设，限制只回答仓储问题，严禁编造数据", 14, DARK_TEXT),
    ("• 真实数据注入（buildDataContext）：智能匹配关键词→查库→注入上下文", 14, DARK_TEXT),
    ("• 上下文对话管理：保留最近10轮历史，localStorage持久化", 14, DARK_TEXT),
    ("• 安全边界：拒绝回答娱乐/金融/政治等无关问题", 14, RED_ACCENT),
])

# 右侧：架构图占位 + 标签
add_image_placeholder(slide, Inches(7.8), Inches(1.5), Inches(4.8), Inches(2.8))
add_tag(slide, Inches(7.8), Inches(4.6), "DeepSeek V3", BLUE_MAIN)
add_tag(slide, Inches(9.8), Inches(4.6), "OkHttp", TEAL)
add_tag(slide, Inches(7.8), Inches(5.1), "System Prompt", BLUE_LIGHT)
add_tag(slide, Inches(9.8), Inches(5.1), "RAG 数据注入", GOLD)

add_multi_text(slide, Inches(7.8), Inches(5.6), Inches(4.8), Inches(1.5), [
    ("对话示例：", 14, WHITE, True),
    ('用户："冷冻鸡排库存多少？"', 13, LIGHT_GRAY),
    ('AI："根据系统数据，冷冻鸡排当前库存\n为200件，分布3个库位..."', 13, TEAL),
])

# ──────────────────────────────────────────────
# 第7页：数据可视化看板（亮点3）
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
             "05  数据可视化看板 —— ECharts 实时监控", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), TEAL)

# 四个图表卡片
charts = [
    ("📈  近7天出入库趋势", "折线图（双线）\n蓝线=入库量  绿线=出库量\n直观展示仓库作业动态", BLUE_MAIN),
    ("📊  货物数量分布", "柱状图\n每种货物的库存数量\n快速了解货物结构", TEAL),
    ("🌡  仓库温湿度监控", "双Y轴折线图\n左轴=温度(℃)  右轴=湿度(%)\n冷链温控可视化", BLUE_LIGHT),
    ("🏭  仓库容量饱和度", "柱状图\n已用库位 vs 总库位\n仓库空间利用率", GOLD),
]

for i, (title, desc, accent) in enumerate(charts):
    x = Inches(0.5) + Inches(3.15) * i
    y = Inches(1.5)
    add_rect(slide, x, y, Inches(2.95), Inches(2.5), BG_CARD)
    add_rect(slide, x, y, Inches(2.95), Pt(4), accent)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.6), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.6), Inches(2.6), Inches(1.8),
                 desc, font_size=12, color=LIGHT_GRAY)

# 底部技术说明
add_multi_text(slide, Inches(0.8), Inches(4.3), Inches(7), Inches(1.2), [
    ("🛠  技术实现", 18, WHITE, True),
    ("• 前端：ECharts 5 图表库，Dashboard 页面懒加载优化首屏性能", 14, DARK_TEXT),
    ("• 后端：4个专属 Dashboard API，对原始数据做 SQL 聚合查询", 14, DARK_TEXT),
    ("• 数据实时刷新，支持响应式布局适配不同屏幕", 14, DARK_TEXT),
])

# 底部图片占位
add_image_placeholder(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.7))

# ──────────────────────────────────────────────
# 第8页：微信小程序
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
             "06  微信小程序 —— 跨平台移动端", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), TEAL)

# 功能介绍
add_multi_text(slide, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.0), [
    ("📱  移动端能力", 22, WHITE, True),
    ("", 6, WHITE),
    ("• 与 Web 端共享同一后端 API，数据实时同步", 16, DARK_TEXT),
    ("• 入库管理 / 出库管理 / 质检管理：支持创建、查看、执行操作", 16, DARK_TEXT),
    ("• 首页数据概览：库存统计 + 快捷操作入口", 16, DARK_TEXT),
    ("• 商品管理 / 批次管理 / 环境监控 / 个人中心", 16, DARK_TEXT),
    ("• 按创建时间倒序排列，本地数据库关联操作人和创建时间", 16, DARK_TEXT),
    ("", 10, WHITE),
    ("🔧  技术特点", 22, WHITE, True),
    ("", 6, WHITE),
    ("• 原生微信小程序框架 + Promise 异步处理", 15, DARK_TEXT),
    ("• wx.request 封装 HTTP 请求，JWT 令牌认证", 15, DARK_TEXT),
    ("• wx.getStorageSync 本地数据持久化，离线也可查看缓存信息", 15, DARK_TEXT),
    ("• 库位可视化选择（按库区分组），条件过滤（储存条件→仓库类型匹配）", 15, DARK_TEXT),
])

# 右侧图片占位
add_image_placeholder(slide, Inches(7.8), Inches(1.5), Inches(4.8), Inches(5.2))

# ──────────────────────────────────────────────
# 第9页：项目亮点总结
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
             "07  项目亮点总结", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), TEAL)

highlights = [
    ("🌟 亮点一", "完整且逻辑闭环的业务流程",
     "从批次创建 → 入库上架 → 库存管理 → 质检把控 → 出库下架 → 报表分析的完整链路。\n"
     "一次入库操作涉及10张表的联动更新，使用 @Transactional 事务保证数据一致性。\n"
     "22张数据库表覆盖仓储管理的所有核心业务场景，三层权限架构清晰分离不同角色职责。",
     GOLD),
    ("🤖 亮点二", "AI 智能助手 —— 真实数据驱动的对话",
     "接入 DeepSeek 大语言模型，通过系统提示词工程限制 AI 职责边界。\n"
     "核心创新：buildDataContext() 方法智能分析用户问题关键词，\n"
     "实时查询数据库并将真实数据注入提示词上下文，AI 基于实际数据回答而非编造。",
     BLUE_MAIN),
    ("📊 亮点三", "数据可视化看板 —— 仓储驾驶舱",
     "4个 ECharts 图表实时展示出入库趋势、货物分布、温湿度监控、仓库容量。\n"
     "Dashboard 页面采用路由懒加载，优化首屏性能。\n"
     "后端通过 SQL 聚合查询为每个图表提供专属数据接口，数据实时准确。",
     TEAL),
    ("📱 亮点四", "Web+小程序双平台覆盖",
     "同一后端同时服务 Web 管理端和微信小程序移动端。\n"
     "小程序支持入库/出库/质检等核心移动操作，库位可视化选择。\n"
     "暗色/亮色主题切换、AES密码加密传输、Excel数据导出等细节完善。",
     BLUE_LIGHT),
]

for i, (badge, title, desc, accent) in enumerate(highlights):
    y = Inches(1.5) + Inches(1.4) * i
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.25), BG_CARD)
    add_rect(slide, Inches(0.8), y, Pt(4), Inches(1.25), accent)
    add_text_box(slide, Inches(1.2), y + Inches(0.08), Inches(1.5), Inches(0.3),
                 badge, font_size=14, color=accent, bold=True)
    add_text_box(slide, Inches(2.6), y + Inches(0.08), Inches(9), Inches(0.3),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.45), Inches(10.5), Inches(0.7),
                 desc, font_size=12, color=LIGHT_GRAY)

# ──────────────────────────────────────────────
# 第10页：系统演示 + 总结
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
             "08  系统演示", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), TEAL)

# 演示要点
add_multi_text(slide, Inches(0.8), Inches(1.5), Inches(6.5), Inches(4.5), [
    ("🖥  现场演示流程", 22, WHITE, True),
    ("", 8, WHITE),
    ("1.  登录系统 → 首页数据概览", 16, DARK_TEXT),
    ("2.  数据看板 → 出入库趋势 / 温湿度 / 仓库容量", 16, DARK_TEXT),
    ("3.  任务中心 → 创建入库任务 → 选择批次 → 执行入库", 16, DARK_TEXT),
    ("4.  任务中心 → 创建出库任务 → 选择库位 → 执行出库", 16, DARK_TEXT),
    ("5.  AI 助手 → 提问「冷冻鸡排库存多少？」 → 展示真实数据回复", 16, DARK_TEXT),
    ("6.  微信小程序 → 入库/出库列表 → 创建操作 → 数据同步", 16, DARK_TEXT),
])

# 右侧图片占位
add_image_placeholder(slide, Inches(7.8), Inches(1.5), Inches(4.8), Inches(5.2))

# ──────────────────────────────────────────────
# 第11页：总结与致谢
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), BLUE_MAIN)

# 中央卡片
add_rect(slide, Inches(2.5), Inches(1.5), Inches(8.5), Inches(4.8), BG_CARD)

add_text_box(slide, Inches(3.0), Inches(2.0), Inches(7.5), Inches(0.8),
             "感谢聆听", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(2.8), Inches(2.5), TEAL)

add_multi_text(slide, Inches(3.0), Inches(3.2), Inches(7.5), Inches(2.5), [
    ("预制菜 WMS 智能仓储管理系统", 22, BLUE_LIGHT, False),
    ("", 10, WHITE),
    ("技术栈：Vue 3 + Spring Boot 3 + MySQL + DeepSeek + 微信小程序", 15, LIGHT_GRAY),
    ("", 6, WHITE),
    ("汇报人：彭健坤    王佳宁    彭双华", 18, WHITE, True),
    ("2026年7月19日", 15, LIGHT_GRAY),
])

add_text_box(slide, Inches(3.0), Inches(6.0), Inches(7.5), Inches(0.5),
             "Q & A", font_size=28, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

# ── 保存 ──
output_path = r"D:\2026summer\from_git\24组演示PPT.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
