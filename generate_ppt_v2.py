#!/usr/bin/env python3
"""
基于中国石油大学汇报答辩通用PPT模板，生成答辩PPT。
保留模板主题（背景图、配色、字体），只替换内容。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── 打开模板 ──
TEMPLATE = '中国石油大学汇报答辩通用ppt.pptx'
prs = Presentation(TEMPLATE)

# 删除模板中所有现有幻灯片
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(qn('r:id'))
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

W = prs.slide_width
H = prs.slide_height

# ── 获取布局 ──
master = prs.slide_masters[0]
layout_cover = master.slide_layouts[0]    # 标题幻灯片
layout_blank = master.slide_layouts[8]    # 空白
layout_end = master.slide_layouts[4]      # 结尾幻灯片
layout_content = master.slide_layouts[6]  # 内容页（仅标题占位）

# ── 颜色定义（与模板协调）──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x1E, 0x56, 0xA0)
BLUE_LIGHT = RGBColor(0x3B, 0x82, 0xF6)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
RED = RGBColor(0xC0, 0x39, 0x2B)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
BG_LIGHT = RGBColor(0xF0, 0xF2, 0xF5)
CARD_BG = RGBColor(0xF8, 0xF9, 0xFA)

# ── 工具函数 ──
def add_rect(slide, left, top, width, height, color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_round_rect(slide, left, top, width, height, color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, size=16, color=DARK, bold=False,
                align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tb

def add_multitext(slide, left, top, width, height, lines):
    """lines: [(text, size, color, bold), ...]"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line[0]
        p.font.size = Pt(line[1])
        p.font.color.rgb = line[2] if len(line) > 2 else DARK
        p.font.bold = line[3] if len(line) > 3 else False
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.LEFT
    return tb

def add_line(slide, left, top, width, color=BLUE, thickness=Pt(3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_img_placeholder(slide, left, top, width, height, text="插入图片"):
    """虚线占位框"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    shape.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    shape.line.width = Pt(1)
    shape.line.dash_style = 2
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"  {text}"
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_tag(slide, left, top, text, color=BLUE):
    shape = add_round_rect(slide, left, top, Inches(1.6), Inches(0.3), color)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER

def add_page_number(slide, num, total=11):
    add_textbox(slide, Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.35),
                f'{num}/{total}', size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
# 第1页：封面（使用模板封面布局）
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(layout_cover)
# 设置标题占位符
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:  # 标题
        ph.text = "预制菜WMS智能仓储管理系统"
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.name = 'Microsoft YaHei'
            run.font.size = Pt(44)
            run.font.color.rgb = WHITE
            run.font.bold = True
    elif ph.placeholder_format.idx == 1:  # 副标题
        ph.text = "Prepared Food Warehouse Management System"
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.name = 'Microsoft YaHei'
            run.font.size = Pt(20)
    elif ph.placeholder_format.idx == 13:  # 文本21
        ph.text = "汇报人：彭健坤    王佳宁    彭双华"
    elif ph.placeholder_format.idx == 15:  # 文本25
        ph.text = "2026年7月19日"

# ═══════════════════════════════════════════
# 第2页：目录
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.6),
            '目  录', size=34, color=DARK, bold=True)
add_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), BLUE)

toc_items = [
    ("01", "项目概述", "背景、目标与系统定位"),
    ("02", "技术架构", "前后端技术栈与系统架构"),
    ("03", "核心业务流程", "出入库、质检、盘点的完整闭环"),
    ("04", "AI 智能助手", "DeepSeek 大模型 + 真实数据注入"),
    ("05", "数据可视化看板", "ECharts 图表实时监控"),
    ("06", "微信小程序", "跨平台移动端适配"),
    ("07", "项目亮点总结", "四大亮点与创新点"),
    ("08", "系统演示", "现场展示 + Q&A"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.4) + Inches(0.7) * i
    add_textbox(slide, Inches(1.0), y, Inches(0.5), Inches(0.35),
                num, size=24, color=BLUE, bold=True)
    add_line(slide, Inches(1.7), y + Inches(0.05), Pt(2), BLUE, Inches(0.25))
    add_textbox(slide, Inches(2.0), y, Inches(3.5), Inches(0.3),
                title, size=18, color=DARK, bold=True)
    add_textbox(slide, Inches(2.0), y + Inches(0.28), Inches(4), Inches(0.25),
                desc, size=11, color=GRAY)

add_img_placeholder(slide, Inches(7.5), Inches(1.4), Inches(5.0), Inches(5.2))
add_page_number(slide, 2)

# ═══════════════════════════════════════════
# 第3页：项目概述
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_textbox(slide, Inches(0.8), Inches(0.35), Inches(6), Inches(0.55),
            '01  项目概述', size=30, color=DARK, bold=True)
add_line(slide, Inches(0.8), Inches(0.9), Inches(1.5), BLUE)

add_multitext(slide, Inches(0.8), Inches(1.3), Inches(6.8), Inches(5.5), [
    (" 项目定位", 20, DARK, True),
    ("为预制菜行业量身打造的仓库管理系统（WMS），覆盖从货物入库、存储、", 14, GRAY),
    ("出库到质检、盘点的全生命周期管理。", 14, GRAY),
    ("", 6, DARK),
    (" 解决的核心问题", 20, DARK, True),
    ("• 传统纸质记录效率低、易出错 → 数字化管理", 14, GRAY),
    ("• 库存数据不透明、难追溯 → 实时库存 + 操作日志", 14, GRAY),
    ("• 多角色协作混乱 → 主管/员工/质检员三级权限", 14, GRAY),
    ("• 缺乏数据支撑决策 → 可视化数据看板", 14, GRAY),
    ("", 6, DARK),
    (" 系统规模", 20, DARK, True),
    ("• 22 张数据库表 | 8 个前端页面 | 6 个小程序页面", 14, GRAY),
    ("• 支持 Web 端 + 微信小程序双平台", 14, GRAY),
    ("• 完整覆盖 入库→存储→出库→质检→盘点 全流程", 14, GRAY),
])

add_img_placeholder(slide, Inches(8.0), Inches(1.3), Inches(4.6), Inches(5.3))
add_page_number(slide, 3)

# ═══════════════════════════════════════════
# 第4页：技术架构
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_textbox(slide, Inches(0.8), Inches(0.35), Inches(6), Inches(0.55),
            '02  技术架构', size=30, color=DARK, bold=True)
add_line(slide, Inches(0.8), Inches(0.9), Inches(1.5), BLUE)

layers = [
    ("  前端层", "Vue 3 + Element Plus + ECharts + Axios",
     "Widescreen (1920x1080), 暗色/亮色主题切换, AES密码加密传输"),
    ("  通信层", "HTTP RESTful API + JWT 令牌认证",
     "Axios 拦截器自动附加 Token, Vite 代理解决跨域, 统一 ResultData 响应格式"),
    ("  后端层", "Spring Boot 3 + MyBatis + Maven",
     "三层架构(Controller→Service→Mapper), 事务管理(@Transactional), 全局异常处理"),
    ("  数据层", "MySQL 8.0 + Redis 6",
     "22张业务表, Redis 缓存用户会话, MyBatis 动态 SQL"),
    ("  AI 层", "DeepSeek 大模型 (API 调用)",
     "OkHttp 发送请求, 系统提示词约束, 真实数据库数据注入上下文"),
]

for i, (title, tech, desc) in enumerate(layers):
    y = Inches(1.2) + Inches(0.95) * i
    add_rect(slide, Inches(0.8), y, Inches(6.8), Inches(0.82), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
    add_line(slide, Inches(0.8), y, Pt(3), BLUE, Inches(0.82))
    add_textbox(slide, Inches(1.1), y + Inches(0.04), Inches(5.5), Inches(0.28),
                title, size=15, color=DARK, bold=True)
    add_textbox(slide, Inches(1.1), y + Inches(0.3), Inches(3.0), Inches(0.22),
                tech, size=12, color=BLUE)
    add_textbox(slide, Inches(1.1), y + Inches(0.52), Inches(5.5), Inches(0.22),
                desc, size=11, color=GRAY)

# 右侧技术栈卡片
add_round_rect(slide, Inches(8.0), Inches(1.2), Inches(4.6), Inches(2.5), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_multitext(slide, Inches(8.3), Inches(1.4), Inches(4.0), Inches(2.2), [
    ("技术栈一览", 17, DARK, True),
    ("", 4, DARK),
    ("前端：Vue 3 · Element Plus · ECharts", 13, GRAY),
    ("后端：Spring Boot 3 · MyBatis", 13, GRAY),
    ("数据库：MySQL · Redis", 13, GRAY),
    ("安全：JWT · BCrypt · AES-256-CBC", 13, BLUE),
    ("AI：DeepSeek · OkHttp", 13, BLUE),
    ("工具：EasyExcel · Vite · Maven", 13, GRAY),
])

add_img_placeholder(slide, Inches(8.0), Inches(4.0), Inches(4.6), Inches(2.7))
add_page_number(slide, 4)

# ═══════════════════════════════════════════
# 第5页：核心业务流程（亮点1）
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.55),
            '03  核心业务流程 —— 完整闭环', size=30, color=DARK, bold=True)
add_line(slide, Inches(0.8), Inches(0.9), Inches(1.5), BLUE)

flow = [
    (" 批次创建", "生产完成后创建批次 记录生产日期、保质期", BLUE),
    (" 入库操作", "选择批次 指定库位 9张表联动更新", TEAL),
    (" 库存管理", "实时库位占用状态 先进先出(FIFO)", BLUE_LIGHT),
    (" 质检管理", "来料检/日常抽检 合格放行 不合格报废", GOLD),
    (" 出库操作", "选择批次 指定库位 库存扣减 日志记录", BLUE),
    (" 报表看板", "数据可视化总览 出入库趋势 温湿度", TEAL),
]

for i, (title, desc, accent) in enumerate(flow):
    x = Inches(0.4) + Inches(2.1) * i
    add_round_rect(slide, x, Inches(1.3), Inches(1.95), Inches(2.1), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
    add_line(slide, x, Inches(1.3), Inches(1.95), accent, Pt(3))
    add_textbox(slide, x + Inches(0.1), Inches(1.45), Inches(1.75), Inches(0.35),
                title, size=15, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), Inches(1.9), Inches(1.75), Inches(1.3),
                desc, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        add_textbox(slide, x + Inches(1.95), Inches(2.0), Inches(0.2), Inches(0.3),
                    '', size=20, color=accent, bold=True, align=PP_ALIGN.CENTER)

# 亮点说明
add_rect(slide, Inches(0.4), Inches(3.7), Inches(12.4), Inches(1.6), RGBColor(0xFF, 0xF8, 0xE1), GOLD)
add_multitext(slide, Inches(0.7), Inches(3.85), Inches(11.8), Inches(1.4), [
    ("  亮点：一次入库操作涉及 10 张数据库表的联动更新，使用 @Transactional 保证事务一致性", 15, DARK, True),
    ("  涉及表：inbound_order_head → inbound_order_detail → work_task → operation_log → inventory →", 12, GRAY),
    ("  location → zone → warehouse → goods → batch（10张表全部通过事务保护，任何一步失败即回滚）", 12, GRAY),
    ("  三层权限架构：主管（全部权限）→ 仓库员工（出入库执行）→ 质检员（质量检查），前端路由守卫 + 后端接口鉴权", 13, GRAY),
])

add_img_placeholder(slide, Inches(0.4), Inches(5.6), Inches(12.4), Inches(1.3), "插入：出入库流程截图")
add_page_number(slide, 5)

# ═══════════════════════════════════════════
# 第6页：AI智能助手（亮点2）
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.55),
            '04  AI 智能助手 —— DeepSeek 大模型接入', size=30, color=DARK, bold=True)
add_line(slide, Inches(0.8), Inches(0.9), Inches(1.5), BLUE)

add_multitext(slide, Inches(0.8), Inches(1.3), Inches(7.0), Inches(5.2), [
    ("  实现原理", 19, DARK, True),
    ("", 4, DARK),
    ("1. 用户输入问题  前端将问题+历史对话发送到后端 /ai/chat", 13, GRAY),
    ("2. 后端分析关键词（库存/货物/仓库/批次） 查询真实数据库", 13, GRAY),
    ("3. 组装消息：系统提示词 + 真实数据上下文 + 历史对话 + 当前问题", 13, GRAY),
    ("4. OkHttp 发送到 DeepSeek API  解析回复  返回前端", 13, GRAY),
    ("", 8, DARK),
    ("  核心技术点", 19, DARK, True),
    ("", 4, DARK),
    ("• 系统提示词工程：定义AI人设，限制只回答仓储问题，严禁编造数据", 13, GRAY),
    ("• 真实数据注入（buildDataContext）：智能匹配关键词 查库 注入上下文", 13, GRAY),
    ("• 上下文对话管理：保留最近10轮历史，localStorage持久化", 13, GRAY),
    ("• 安全边界：拒绝回答娱乐/金融/政治等无关问题", 13, RED),
])

add_img_placeholder(slide, Inches(8.2), Inches(1.3), Inches(4.5), Inches(2.7), "插入：AI对话界面截图")

add_tag(slide, Inches(8.2), Inches(4.2), "DeepSeek V3", BLUE)
add_tag(slide, Inches(10.0), Inches(4.2), "OkHttp", TEAL)
add_tag(slide, Inches(8.2), Inches(4.65), "System Prompt", BLUE_LIGHT)
add_tag(slide, Inches(10.0), Inches(4.65), "RAG 数据注入", GOLD)

add_multitext(slide, Inches(8.2), Inches(5.2), Inches(4.5), Inches(1.5), [
    ("对话示例：", 13, DARK, True),
    ('用户："冷冻鸡排库存多少？"', 12, GRAY),
    ('AI："根据系统数据，冷冻鸡排当前库存 为200件，分布3个库位..."', 12, TEAL),
])
add_page_number(slide, 6)

# ═══════════════════════════════════════════
# 第7页：数据可视化看板（亮点3）
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.55),
            '05  数据可视化看板 —— ECharts 实时监控', size=30, color=DARK, bold=True)
add_line(slide, Inches(0.8), Inches(0.9), Inches(1.5), BLUE)

charts = [
    ("  近7天出入库趋势", "折线图（双线） 蓝线=入库量 绿线=出库量  直观展示仓库作业动态", BLUE),
    ("  货物数量分布", "柱状图  每种货物的库存数量  快速了解货物结构", TEAL),
    ("  仓库温湿度监控", "双Y轴折线图  左轴=温度( ) 右轴=湿度(%)  冷链温控可视化", BLUE_LIGHT),
    ("  仓库容量饱和度", "柱状图  已用库位 vs 总库位  仓库空间利用率", GOLD),
]

for i, (title, desc, accent) in enumerate(charts):
    x = Inches(0.4) + Inches(3.15) * i
    add_round_rect(slide, x, Inches(1.3), Inches(3.0), Inches(2.3), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
    add_line(slide, x, Inches(1.3), Inches(3.0), accent, Pt(3))
    add_textbox(slide, x + Inches(0.15), Inches(1.45), Inches(2.7), Inches(0.35),
                title, size=14, color=DARK, bold=True)
    add_textbox(slide, x + Inches(0.15), Inches(1.9), Inches(2.7), Inches(1.5),
                desc, size=11, color=GRAY)

add_multitext(slide, Inches(0.8), Inches(3.9), Inches(7), Inches(1.2), [
    ("  技术实现", 17, DARK, True),
    ("• 前端：ECharts 5 图表库，Dashboard 页面懒加载优化首屏性能", 13, GRAY),
    ("• 后端：4个专属 Dashboard API，对原始数据做 SQL 聚合查询", 13, GRAY),
    ("• 数据实时刷新，支持响应式布局适配不同屏幕", 13, GRAY),
])

add_img_placeholder(slide, Inches(0.4), Inches(5.2), Inches(12.4), Inches(1.7), "插入：数据看板全景截图")
add_page_number(slide, 7)

# ═══════════════════════════════════════════
# 第8页：微信小程序
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.55),
            '06  微信小程序 —— 跨平台移动端', size=30, color=DARK, bold=True)
add_line(slide, Inches(0.8), Inches(0.9), Inches(1.5), BLUE)

add_multitext(slide, Inches(0.8), Inches(1.3), Inches(7.0), Inches(5.2), [
    ("  移动端能力", 19, DARK, True),
    ("", 4, DARK),
    ("• 与 Web 端共享同一后端 API，数据实时同步", 14, GRAY),
    ("• 入库管理 / 出库管理 / 质检管理：支持创建、查看、执行操作", 14, GRAY),
    ("• 首页数据概览：库存统计 + 快捷操作入口", 14, GRAY),
    ("• 商品管理 / 批次管理 / 环境监控 / 个人中心", 14, GRAY),
    ("• 按创建时间倒序排列，本地数据库关联操作人和创建时间", 14, GRAY),
    ("", 8, DARK),
    ("  技术特点", 19, DARK, True),
    ("", 4, DARK),
    ("• 原生微信小程序框架 + Promise 异步处理", 14, GRAY),
    ("• wx.request 封装 HTTP 请求，JWT 令牌认证", 14, GRAY),
    ("• wx.getStorageSync 本地数据持久化，离线也可查看缓存信息", 14, GRAY),
    ("• 库位可视化选择（按库区分组），条件过滤（储存条件 仓库类型匹配）", 14, GRAY),
])

add_img_placeholder(slide, Inches(8.2), Inches(1.3), Inches(4.5), Inches(5.3), "插入：小程序截图合集")
add_page_number(slide, 8)

# ═══════════════════════════════════════════
# 第9页：项目亮点总结
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_textbox(slide, Inches(0.8), Inches(0.35), Inches(6), Inches(0.55),
            '07  项目亮点总结', size=30, color=DARK, bold=True)
add_line(slide, Inches(0.8), Inches(0.9), Inches(1.5), BLUE)

highlights = [
    ("  亮点一", "完整且逻辑闭环的业务流程",
     "从批次创建  入库上架  库存管理  质检把控  出库下架  报表分析的完整链路。"
     "一次入库操作涉及10张表的联动更新，使用 @Transactional 事务保证数据一致性。"
     "22张数据库表覆盖仓储管理的所有核心业务场景，三层权限架构清晰分离不同角色职责。",
     GOLD),
    ("  亮点二", "AI 智能助手 —— 真实数据驱动的对话",
     "接入 DeepSeek 大语言模型，通过系统提示词工程限制 AI 职责边界。"
     "核心创新：buildDataContext() 方法智能分析用户问题关键词，"
     "实时查询数据库并将真实数据注入提示词上下文，AI 基于实际数据回答而非编造。",
     BLUE),
    ("  亮点三", "数据可视化看板 —— 仓储驾驶舱",
     "4个 ECharts 图表实时展示出入库趋势、货物分布、温湿度监控、仓库容量。"
     "Dashboard 页面采用路由懒加载，优化首屏性能。"
     "后端通过 SQL 聚合查询为每个图表提供专属数据接口，数据实时准确。",
     TEAL),
    ("  亮点四", "Web+小程序双平台覆盖",
     "同一后端同时服务 Web 管理端和微信小程序移动端。"
     "小程序支持入库/出库/质检等核心移动操作，库位可视化选择。"
     "暗色/亮色主题切换、AES密码加密传输、Excel数据导出等细节完善。",
     BLUE_LIGHT),
]

for i, (badge, title, desc, accent) in enumerate(highlights):
    y = Inches(1.2) + Inches(1.45) * i
    add_round_rect(slide, Inches(0.8), y, Inches(11.8), Inches(1.3), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
    add_line(slide, Inches(0.8), y, Pt(3), accent, Inches(1.3))
    add_textbox(slide, Inches(1.2), y + Inches(0.06), Inches(1.5), Inches(0.25),
                badge, size=13, color=accent, bold=True)
    add_textbox(slide, Inches(2.6), y + Inches(0.06), Inches(9.5), Inches(0.25),
                title, size=16, color=DARK, bold=True)
    add_textbox(slide, Inches(1.2), y + Inches(0.42), Inches(10.8), Inches(0.75),
                desc, size=11, color=GRAY)
add_page_number(slide, 9)

# ═══════════════════════════════════════════
# 第10页：系统演示
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_textbox(slide, Inches(0.8), Inches(0.35), Inches(6), Inches(0.55),
            '08  系统演示', size=30, color=DARK, bold=True)
add_line(slide, Inches(0.8), Inches(0.9), Inches(1.5), BLUE)

add_multitext(slide, Inches(0.8), Inches(1.3), Inches(7.0), Inches(5.0), [
    ("  现场演示流程", 19, DARK, True),
    ("", 6, DARK),
    ("1.  登录系统  首页数据概览", 15, GRAY),
    ("2.  数据看板  出入库趋势 / 温湿度 / 仓库容量", 15, GRAY),
    ("3.  任务中心  创建入库任务  选择批次  执行入库", 15, GRAY),
    ("4.  任务中心  创建出库任务  选择库位  执行出库", 15, GRAY),
    ("5.  AI 助手  提问  展示真实数据回复", 15, GRAY),
    ("6.  微信小程序  入库/出库列表  创建操作  数据同步", 15, GRAY),
])

add_img_placeholder(slide, Inches(8.2), Inches(1.3), Inches(4.5), Inches(5.8), "插入：系统演示准备")
add_page_number(slide, 10)

# ═══════════════════════════════════════════
# 第11页：总结与致谢（使用模板结尾布局）
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(layout_end)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:  # 标题
        ph.text = "感谢聆听"
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.name = 'Microsoft YaHei'
            run.font.size = Pt(48)
            run.font.color.rgb = WHITE
            run.font.bold = True
    elif ph.placeholder_format.idx == 1:  # 副标题
        ph.text = "Q & A"
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.name = 'Microsoft YaHei'
            run.font.size = Pt(28)

# ── 保存 ──
OUTPUT = r'D:\2026summer\from_git\24组演示PPT_new.pptx'
prs.save(OUTPUT)
print(f'Done! Saved to: {OUTPUT}')
print(f'Total slides: {len(prs.slides)}')
